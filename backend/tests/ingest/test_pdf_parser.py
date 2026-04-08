"""
test_pdf_parser.py — 覆蓋率目標 80%

PDF：mock GemmaClient（不呼叫真實 Ollama）
PPT / DOCX：純文字提取，不需 mock
"""
import io
import os
import tempfile
from unittest.mock import AsyncMock, MagicMock, patch

import fitz
import pytest
from docx import Document
from pptx import Presentation

from backend.ingest.base import SourceType
from backend.ingest.pdf_parser import process, _render_page


# ── Helper：建立測試檔案 ──────────────────────────────────────────────────────

def _make_pdf(pages: int = 2) -> str:
    """建立暫存 PDF（每頁一個頁碼數字），回傳路徑"""
    tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
    tmp.close()
    doc = fitz.open()
    for i in range(pages):
        page = doc.new_page()
        page.insert_text((50, 100), f"Page {i + 1}")
    doc.save(tmp.name)
    doc.close()
    return tmp.name


def _make_pptx(slides: list[str]) -> str:
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp.close()
    prs = Presentation()
    for text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        txBox.text_frame.text = text
    prs.save(tmp.name)
    return tmp.name


def _make_docx(sections: list[tuple[str, str]]) -> str:
    tmp = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
    tmp.close()
    doc = Document()
    for heading, body in sections:
        doc.add_heading(heading, level=1)
        doc.add_paragraph(body)
    doc.save(tmp.name)
    return tmp.name


# ── PDF 測試（mock GemmaClient）──────────────────────────────────────────────

def _mock_gemma(return_text: str = "Extracted text from page"):
    """回傳 patch context，mock GemmaClient.generate"""
    mock_client = MagicMock()
    mock_client.generate = AsyncMock(return_value=return_text)
    return patch("backend.ingest.pdf_parser.GemmaClient", return_value=mock_client)


async def test_pdf_calls_gemma_per_page():
    """每一頁都應呼叫一次 GemmaClient.generate"""
    path = _make_pdf(pages=3)
    try:
        with _mock_gemma("page content") as MockGemma:
            chunks = await process(path)
        instance = MockGemma.return_value
        assert instance.generate.call_count == 3
        assert len(chunks) == 3
    finally:
        os.unlink(path)


async def test_pdf_chunk_source_type():
    path = _make_pdf(pages=1)
    try:
        with _mock_gemma("some text"):
            chunks = await process(path)
        assert chunks[0].source_type == SourceType.PDF
    finally:
        os.unlink(path)


async def test_pdf_chunk_metadata():
    path = _make_pdf(pages=2)
    try:
        with _mock_gemma("text"):
            chunks = await process(path)
        assert chunks[0].metadata["page"] == 1
        assert chunks[1].metadata["page"] == 2
        assert chunks[0].metadata["total_pages"] == 2
        assert "filename" in chunks[0].metadata
    finally:
        os.unlink(path)


async def test_pdf_source_id_format():
    path = _make_pdf(pages=1)
    try:
        with _mock_gemma("text"):
            chunks = await process(path)
        assert "::p1" in chunks[0].source_id
    finally:
        os.unlink(path)


async def test_pdf_empty_gemma_response_skips_chunk():
    """Gemma 回傳空字串時不應產生 chunk"""
    path = _make_pdf(pages=2)
    try:
        mock_client = MagicMock()
        mock_client.generate = AsyncMock(side_effect=["", "real content"])
        with patch("backend.ingest.pdf_parser.GemmaClient", return_value=mock_client):
            chunks = await process(path)
        assert len(chunks) == 1
        assert chunks[0].content == "real content"
    finally:
        os.unlink(path)


async def test_pdf_gemma_error_skips_page_continues():
    """單頁 Gemma 失敗不應中斷整份 PDF"""
    from backend.engine.gemma_client import GemmaEdgeUnavailable
    path = _make_pdf(pages=3)
    try:
        mock_client = MagicMock()
        mock_client.generate = AsyncMock(
            side_effect=[GemmaEdgeUnavailable("down"), "page 2 text", "page 3 text"]
        )
        with patch("backend.ingest.pdf_parser.GemmaClient", return_value=mock_client):
            chunks = await process(path)
        assert len(chunks) == 2
    finally:
        os.unlink(path)


async def test_pdf_images_passed_to_gemma():
    """每次呼叫 generate 應傳入 images 參數"""
    path = _make_pdf(pages=1)
    try:
        mock_client = MagicMock()
        mock_client.generate = AsyncMock(return_value="text")
        with patch("backend.ingest.pdf_parser.GemmaClient", return_value=mock_client):
            await process(path)
        call_kwargs = mock_client.generate.call_args
        assert "images" in call_kwargs.kwargs
        assert len(call_kwargs.kwargs["images"]) == 1
        assert call_kwargs.kwargs["mode"] == "edge"
    finally:
        os.unlink(path)


# ── _render_page 單元測試 ────────────────────────────────────────────────────

def test_render_page_returns_png_bytes():
    doc = fitz.open()
    page = doc.new_page()
    result = _render_page(page)
    assert isinstance(result, bytes)
    assert result[:4] == b"\x89PNG"  # PNG magic bytes


def test_render_page_size_reasonable():
    """渲染後的 PNG 應大於 1KB（確保有實際內容）"""
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 100), "Hello World")
    result = _render_page(page)
    assert len(result) > 1024


# ── PPT 測試（純文字，不需 mock）─────────────────────────────────────────────

async def test_pptx_basic():
    path = _make_pptx(["Slide one content", "Slide two content"])
    try:
        chunks = await process(path)
        assert len(chunks) == 2
        assert chunks[0].source_type == SourceType.PPT
        assert chunks[0].metadata["slide"] == 1
        assert chunks[1].metadata["slide"] == 2
    finally:
        os.unlink(path)


async def test_pptx_content_extracted():
    path = _make_pptx(["derivative is the slope of tangent"])
    try:
        chunks = await process(path)
        assert "derivative" in chunks[0].content.lower()
    finally:
        os.unlink(path)


async def test_pptx_empty_slides_skipped():
    path = _make_pptx(["has content", ""])
    try:
        chunks = await process(path)
        assert len(chunks) == 1
    finally:
        os.unlink(path)


async def test_pptx_source_id_format():
    path = _make_pptx(["content"])
    try:
        chunks = await process(path)
        assert "::slide1" in chunks[0].source_id
    finally:
        os.unlink(path)


# ── DOCX 測試（純文字，不需 mock）────────────────────────────────────────────

async def test_docx_basic():
    path = _make_docx([("Chapter 1", "content here"), ("Chapter 2", "more content")])
    try:
        chunks = await process(path)
        assert len(chunks) >= 1
        assert all(c.source_type == SourceType.WORD for c in chunks)
    finally:
        os.unlink(path)


async def test_docx_content_extracted():
    path = _make_docx([("Title", "this is the body paragraph")])
    try:
        chunks = await process(path)
        assert any("body paragraph" in c.content for c in chunks)
    finally:
        os.unlink(path)


async def test_docx_metadata_has_section():
    path = _make_docx([("Chapter 1", "content")])
    try:
        chunks = await process(path)
        assert "section" in chunks[0].metadata
        assert "filename" in chunks[0].metadata
    finally:
        os.unlink(path)


# ── 錯誤處理 ──────────────────────────────────────────────────────────────────

async def test_file_not_found():
    with pytest.raises(FileNotFoundError):
        await process("/nonexistent/file.pdf")


async def test_unsupported_format():
    tmp = tempfile.NamedTemporaryFile(suffix=".txt", delete=False)
    tmp.write(b"text")
    tmp.close()
    try:
        with pytest.raises(ValueError, match="不支援的格式"):
            await process(tmp.name)
    finally:
        os.unlink(tmp.name)
