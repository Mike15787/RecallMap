"""
PDF / PPT / Word 文字萃取模組
RUNTIME: edge（PDF 整頁渲染 → Gemma 4 多模態；PPT / DOCX 純文字提取）

PDF 策略（與閉源模型相同路線）：
  每頁 get_pixmap() → PNG bytes → GemmaClient.generate(images=[...], mode="edge")
  模型同時理解文字、圖表、公式，無需分兩路處理。
  需要 edge model（Ollama / llama.cpp / vLLM）在線。

PPT / DOCX 策略：
  純文字提取（不依賴模型），圖片 shape 暫不處理（P2 排程）。
"""
import logging
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

from .base import ChunkList, DocumentChunk, SourceType
from backend.engine.gemma_client import GemmaClient

logger = logging.getLogger(__name__)

# PDF 頁面渲染比例（2x ≈ 144 DPI，足夠讀清楚文字，payload 合理）
RENDER_SCALE = 2.0

_PDF_PROMPT = (
    "請仔細分析這張 PDF 頁面的完整內容。"
    "提取所有文字（標題、正文、條列、表格）。"
    "若有圖表、示意圖或公式，用文字簡要描述其內容與意義。"
    "保留原始結構，以純文字輸出，不要加任何額外說明。"
    "若辨識不確定，標記 [?]。"
)


async def process(source: str | Path) -> ChunkList:
    """
    單一公開入口。
    根據副檔名自動選擇解析策略，回傳 DocumentChunk 列表。
    """
    path = Path(source)
    if not path.exists():
        raise FileNotFoundError(f"找不到檔案：{path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return await _parse_pdf(path)
    elif suffix in (".ppt", ".pptx"):
        return _parse_pptx(path)
    elif suffix in (".doc", ".docx"):
        return _parse_docx(path)
    else:
        raise ValueError(f"不支援的格式：{suffix}，請使用 PDF / PPT / PPTX / DOC / DOCX")


# ── PDF：整頁渲染 → Gemma ─────────────────────────────────────────────────────

async def _parse_pdf(path: Path) -> ChunkList:
    """每頁渲染成 PNG，送 Gemma edge model 進行視覺理解"""
    client = GemmaClient()
    chunks: ChunkList = []

    with fitz.open(path) as doc:
        total = len(doc)
        for page_num, page in enumerate(doc, start=1):
            image_bytes = _render_page(page)
            try:
                text = await client.generate(
                    prompt=_PDF_PROMPT,
                    images=[image_bytes],
                    mode="edge",
                )
                if text.strip():
                    chunks.append(
                        DocumentChunk(
                            content=text.strip(),
                            source_type=SourceType.PDF,
                            source_id=f"{path.stem}::p{page_num}",
                            metadata={
                                "filename": path.name,
                                "page": page_num,
                                "total_pages": total,
                            },
                        )
                    )
            except Exception as e:
                logger.warning(f"[pdf_parser] 第 {page_num} 頁處理失敗，跳過：{e}")

    return chunks


def _render_page(page: fitz.Page) -> bytes:
    """將 PDF 頁面渲染成 PNG bytes（RENDER_SCALE 倍解析度）"""
    matrix = fitz.Matrix(RENDER_SCALE, RENDER_SCALE)
    pix = page.get_pixmap(matrix=matrix)
    return pix.tobytes("png")


# ── PPT：純文字提取（圖片 shape 暫不處理）────────────────────────────────────

def _parse_pptx(path: Path) -> ChunkList:
    """
    從 PPTX 提取各投影片的文字 shape 內容。
    圖片 shape / SmartArt / OLE 暫不處理（P2 排程）。
    """
    chunks: ChunkList = []
    prs = Presentation(path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if not texts:
            continue
        chunks.append(
            DocumentChunk(
                content="\n".join(texts),
                source_type=SourceType.PPT,
                source_id=f"{path.stem}::slide{slide_num}",
                metadata={"filename": path.name, "slide": slide_num},
            )
        )
    return chunks


# ── DOCX：依 Heading 切 chunk ─────────────────────────────────────────────────

def _parse_docx(path: Path) -> ChunkList:
    """每個 Heading 切一個 chunk；沒有 Heading 則整份為一個 chunk"""
    doc = Document(path)
    chunks: ChunkList = []
    current_heading = "文件開頭"
    current_lines: list[str] = []
    section_num = 1

    def _flush(heading: str, lines: list[str]) -> None:
        nonlocal section_num
        text = "\n".join(lines).strip()
        if text:
            chunks.append(
                DocumentChunk(
                    content=text,
                    source_type=SourceType.WORD,
                    source_id=f"{path.stem}::s{section_num}",
                    metadata={"filename": path.name, "section": heading},
                )
            )
            section_num += 1

    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            _flush(current_heading, current_lines)
            current_heading = para.text.strip() or current_heading
            current_lines = []
        else:
            if para.text.strip():
                current_lines.append(para.text.strip())

    _flush(current_heading, current_lines)
    return chunks
